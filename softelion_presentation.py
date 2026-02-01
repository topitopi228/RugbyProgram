"""
SoftElion Company Presentation Generator
Creates a modern PowerPoint presentation for IT outsourcing company
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def add_gradient_background(slide, color1=(16, 185, 129), color2=(59, 130, 246)):
    """Add gradient-like background using shapes"""
    left = top = 0
    width = Inches(10)
    height = Inches(7.5)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color1)
    shape.line.color.rgb = RGBColor(*color1)
    
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.75), width, Inches(3.75))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(*color2)
    overlay.fill.transparency = 0.7
    overlay.line.fill.background()
    
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    slide.shapes._spTree.remove(overlay._element)
    slide.shapes._spTree.insert(3, overlay._element)

def add_decorative_elements(slide):
    """Add modern decorative elements"""
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0.5), Inches(1.5), Inches(1.5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(251, 191, 36)
    circle1.fill.transparency = 0.85
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(5.5), Inches(2), Inches(2))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(167, 139, 250)
    circle2.fill.transparency = 0.9
    circle2.line.fill.background()

def add_text_with_style(slide, text, x, y, width, height, size=24, bold=False, color=(255,255,255), align=PP_ALIGN.LEFT):
    """Helper function to add styled text"""
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(size)
    text_frame.paragraphs[0].font.bold = bold
    text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)
    text_frame.paragraphs[0].alignment = align
    return text_frame

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (16, 185, 129), (59, 130, 246))
    add_decorative_elements(slide)
    
    add_text_with_style(slide, "SoftElion", 1, 2, 8, 1.5, 72, True, (255,255,255), PP_ALIGN.CENTER)
    add_text_with_style(slide, "IT Outsourcing Company", 1, 3.5, 8, 1, 32, False, (229,231,235), PP_ALIGN.CENTER)
    add_text_with_style(slide, "Custom Software Development Services", 1, 4.5, 8, 1, 24, False, (209,213,219), PP_ALIGN.CENTER)
    
    # Slide 2: Who We Are
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (30, 41, 59), (16, 185, 129))
    
    add_text_with_style(slide, "Хто ми такі", 0.5, 0.5, 9, 1, 48, True, (255,255,255), PP_ALIGN.CENTER)
    
    info = [
        ("🚀 Заснована у 2020", "4+ роки досвіду на ринку"),
        ("🌍 Міжнародна компанія", "Працюємо з клієнтами по всьому світу"),
        ("⭐ Рейтинг 4.9/5", "Від 150+ задоволених клієнтів"),
        ("✅ 500+ проектів", "Успішно доставлено")
    ]
    
    y = 2
    for title, desc in info:
        tf = add_text_with_style(slide, title, 1.5, y, 7, 0.4, 24, True, (251,191,36))
        p = tf.add_paragraph()
        p.text = f"     {desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(229, 231, 235)
        y += 1.2
    
    # Slide 3: What We Do
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (59, 130, 246), (167, 139, 250))
    
    add_text_with_style(slide, "Що ми робимо", 0.5, 0.5, 9, 1, 48, True, (255,255,255), PP_ALIGN.CENTER)
    
    services = [
        "💻 Custom Software Development",
        "🌐 Web Development",
        "📱 Mobile App Development",
        "🤖 AI & Machine Learning",
        "⚙️ DevOps Services",
        "👥 Dedicated Teams",
        "🔧 Legacy Modernization",
        "✔️ Quality Assurance",
        "🔌 API Development"
    ]
    
    for i, service in enumerate(services):
        x = 1.5 if i % 2 == 0 else 5.5
        y = 2 + (i // 2) * 0.7
        add_text_with_style(slide, service, x, y, 3.5, 0.6, 18, False, (255,255,255))
    
    # Slide 4: Development Process
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (16, 185, 129), (30, 41, 59))
    
    add_text_with_style(slide, "Custom Software Development", 0.5, 0.3, 9, 0.8, 40, True, (255,255,255), PP_ALIGN.CENTER)
    
    process = [
        ("1️⃣ Аналіз вимог", "Детальне вивчення бізнес-потреб"),
        ("2️⃣ UI/UX дизайн", "Створення інтуїтивних інтерфейсів"),
        ("3️⃣ Розробка", "Agile/Scrum методологія"),
        ("4️⃣ Тестування", "Автоматизоване та мануальне"),
        ("5️⃣ Деплоймент", "CI/CD pipeline"),
        ("6️⃣ Підтримка", "24/7 моніторинг")
    ]
    
    y = 1.5
    for step, desc in process:
        tf = add_text_with_style(slide, step, 1, y, 8, 0.4, 20, True, (251,191,36))
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(229, 231, 235)
        y += 0.9
    
    # Slide 5: Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (167, 139, 250), (59, 130, 246))
    
    add_text_with_style(slide, "Технологічний стек", 0.5, 0.3, 9, 0.8, 44, True, (255,255,255), PP_ALIGN.CENTER)
    
    tech = [
        ("Frontend", "React, Angular, Vue.js, Next.js, TypeScript"),
        ("Backend", "Node.js, Python, Java, .NET, PHP, Go"),
        ("Mobile", "React Native, Flutter, Swift, Kotlin"),
        ("Databases", "PostgreSQL, MongoDB, Redis, MySQL"),
        ("Cloud & DevOps", "AWS, Azure, GCP, Docker, Kubernetes"),
        ("AI & ML", "TensorFlow, PyTorch, OpenAI, Scikit-learn")
    ]
    
    y = 1.5
    for category, techs in tech:
        add_text_with_style(slide, f"▶ {category}", 1, y, 8, 0.4, 20, True, (251,191,36))
        add_text_with_style(slide, techs, 1.5, y + 0.35, 7.5, 0.4, 16, False, (229,231,235))
        y += 0.95
    
    # Slide 6: Web & Mobile
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (251, 146, 60), (217, 70, 239))
    
    add_text_with_style(slide, "Web & Mobile Development", 0.5, 0.3, 9, 0.8, 42, True, (255,255,255), PP_ALIGN.CENTER)
    
    tf1 = add_text_with_style(slide, "🌐 Web Development", 0.5, 1.5, 4.5, 0.5, 24, True, (251,191,36))
    web_items = ["• Progressive Web Apps", "• E-commerce платформи", "• SaaS рішення", "• Real-time додатки", "• SEO оптимізація"]
    for item in web_items:
        p = tf1.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(229, 231, 235)
    
    tf2 = add_text_with_style(slide, "📱 Mobile Development", 5, 1.5, 4.5, 0.5, 24, True, (251,191,36))
    mobile_items = ["• iOS & Android", "• Cross-platform", "• Нативна розробка", "• Push-повідомлення", "• App Store оптимізація"]
    for item in mobile_items:
        p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(229, 231, 235)
    
    # Slide 7: AI Solutions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (99, 102, 241), (168, 85, 247))
    
    add_text_with_style(slide, "AI & Machine Learning", 0.5, 0.3, 9, 0.8, 42, True, (255,255,255), PP_ALIGN.CENTER)
    
    ai = [
        ("🤖 Чат-боти", "GPT інтеграція, NLP"),
        ("📊 Предиктивна аналітика", "Прогнозування трендів"),
        ("👁️ Computer Vision", "Розпізнавання об'єктів"),
        ("🗣️ NLP", "Аналіз тексту, переклад"),
        ("🎯 Рекомендації", "Персоналізація контенту"),
        ("🔄 Автоматизація", "RPA, інтелектуальна автоматизація")
    ]
    
    y = 1.5
    for title, desc in ai:
        tf = add_text_with_style(slide, title, 1, y, 8, 0.4, 18, True, (251,191,36))
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(229, 231, 235)
        y += 0.9
    
    # Slide 8: DevOps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (16, 185, 129), (59, 130, 246))
    
    add_text_with_style(slide, "DevOps & Cloud Services", 0.5, 0.3, 9, 0.8, 42, True, (255,255,255), PP_ALIGN.CENTER)
    
    devops = [
        ("☁️ Cloud Migration", "AWS, Azure, GCP"),
        ("🔄 CI/CD", "Jenkins, GitLab CI"),
        ("🐳 Containerization", "Docker, Kubernetes"),
        ("📈 Monitoring", "Prometheus, Grafana"),
        ("🔒 Security", "Infrastructure as Code"),
        ("⚡ Performance", "Auto-scaling, CDN")
    ]
    
    y = 1.5
    for title, desc in devops:
        tf = add_text_with_style(slide, title, 1, y, 8, 0.4, 18, True, (251,191,36))
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(229, 231, 235)
        y += 0.9
    
    # Slide 9: Dedicated Teams
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (217, 70, 239), (99, 102, 241))
    
    add_text_with_style(slide, "Dedicated Teams", 0.5, 0.3, 9, 0.8, 44, True, (255,255,255), PP_ALIGN.CENTER)
    
    benefits = [
        "✅ Повний контроль над командою",
        "✅ Гнучке масштабування",
        "✅ Прозора комунікація",
        "✅ Економія до 60% бюджету",
        "✅ Швидкий старт (2-4 тижні)",
        "✅ Виділені експерти"
    ]
    
    y = 2
    for benefit in benefits:
        add_text_with_style(slide, benefit, 2, y, 6, 0.5, 20, False, (255,255,255))
        y += 0.7
    
    add_text_with_style(slide, "PM • Tech Lead • Developers • QA • DevOps • UI/UX", 1, 6, 8, 1, 16, False, (251,191,36), PP_ALIGN.CENTER)
    
    # Slide 10: Why Choose Us
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (30, 41, 59), (16, 185, 129))
    
    add_text_with_style(slide, "Чому обирають SoftElion?", 0.5, 0.3, 9, 0.8, 44, True, (255,255,255), PP_ALIGN.CENTER)
    
    reasons = [
        ("🏆 500+ проектів", "Доведена експертиза"),
        ("⚡ Time-to-Market", "Швидка доставка MVP"),
        ("💰 Оптимальна ціна", "До 60% економії"),
        ("🛡️ Гарантія якості", "ISO стандарти"),
        ("🌍 24/7 підтримка", "Різні часові зони"),
        ("🤝 Agile", "Щоденні звіти")
    ]
    
    y = 1.5
    for title, desc in reasons:
        tf = add_text_with_style(slide, title, 1, y, 8, 0.4, 20, True, (251,191,36))
        p = tf.add_paragraph()
        p.text = f"    {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(229, 231, 235)
        y += 0.9
    
    # Slide 11: Contact
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (16, 185, 129), (59, 130, 246))
    add_decorative_elements(slide)
    
    add_text_with_style(slide, "Let's Build Together", 1, 2, 8, 1, 48, True, (255,255,255), PP_ALIGN.CENTER)
    add_text_with_style(slide, "🌐 www.softelion.com", 1, 3.5, 8, 0.7, 28, False, (251,191,36), PP_ALIGN.CENTER)
    add_text_with_style(slide, "📧 contact@softelion.com", 1, 4.3, 8, 0.7, 24, False, (229,231,235), PP_ALIGN.CENTER)
    add_text_with_style(slide, "Transform Your Ideas Into Reality", 1, 5.5, 8, 0.7, 20, False, (209,213,219), PP_ALIGN.CENTER)
    
    # Save presentation
    prs.save('SoftElion_Presentation.pptx')
    print("Presentation successfully created: SoftElion_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
